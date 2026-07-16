import fitz  # PyMuPDF
from docx import Document as DocxDocument
from pptx import Presentation
from langchain_core.documents import Document


class DocumentLoader:
    """Loads different document types into LangChain Documents."""

    @staticmethod
    def load_pdf(file):
        documents = []

        pdf = fitz.open(stream=file.read(), filetype="pdf")

        for page_num, page in enumerate(pdf):
            text = page.get_text()

            if text.strip():
                documents.append(
                    Document(
                        page_content=text,
                        metadata={
                            "source": file.name,
                            "page": page_num + 1,
                            "type": "pdf"
                        }
                    )
                )

        return documents

    @staticmethod
    def load_docx(file):
        doc = DocxDocument(file)

        text = "\n".join(
            para.text for para in doc.paragraphs if para.text.strip()
        )

        return [
            Document(
                page_content=text,
                metadata={
                    "source": file.name,
                    "page": 1,
                    "type": "docx"
                }
            )
        ]

    @staticmethod
    def load_pptx(file):
        presentation = Presentation(file)

        documents = []

        for slide_number, slide in enumerate(presentation.slides):

            slide_text = []

            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    if shape.text.strip():
                        slide_text.append(shape.text)

            documents.append(
                Document(
                    page_content="\n".join(slide_text),
                    metadata={
                        "source": file.name,
                        "page": slide_number + 1,
                        "type": "pptx"
                    }
                )
            )

        return documents

    @staticmethod
    def load_txt(file):
        text = file.read().decode("utf-8")

        return [
            Document(
                page_content=text,
                metadata={
                    "source": file.name,
                    "page": 1,
                    "type": "txt"
                }
            )
        ]

    @staticmethod
    def load_documents(uploaded_files):
        all_documents = []

        for file in uploaded_files:

            extension = file.name.split(".")[-1].lower()

            if extension == "pdf":
                all_documents.extend(DocumentLoader.load_pdf(file))

            elif extension == "docx":
                all_documents.extend(DocumentLoader.load_docx(file))

            elif extension == "pptx":
                all_documents.extend(DocumentLoader.load_pptx(file))

            elif extension == "txt":
                all_documents.extend(DocumentLoader.load_txt(file))

        return all_documents